#!/usr/bin/env python3
"""Generate Dealality AI Demand Positioning owner pitch deck (dark product UI PPTX)."""

from __future__ import annotations

import glob
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]

# Dealality brand tokens (from ai-visibility-shared.css)
C_BG = RGBColor(0x08, 0x0F, 0x25)
C_SURFACE = RGBColor(0x0C, 0x14, 0x2C)
C_SURFACE2 = RGBColor(0x21, 0x2C, 0x4D)
C_BORDER = RGBColor(0x37, 0x44, 0x6B)
C_TEXT = RGBColor(0xAE, 0xB9, 0xE1)
C_TEXT_BRIGHT = RGBColor(0xD9, 0xE1, 0xFA)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_ACCENT = RGBColor(0x6C, 0x72, 0xFF)
C_CYAN = RGBColor(0x57, 0xC3, 0xFF)
C_GOLD = RGBColor(0xFD, 0xB5, 0x2A)
C_GREEN = RGBColor(0x14, 0xCA, 0x74)
C_RED = RGBColor(0xDC, 0x2B, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Segoe UI"
DATE_LABEL = "August 2026"
OUTPUT_FILENAME = "Dealality - AI Demand Positioning Deck (Owners) - Dark Product UI.pptx"


def set_slide_bg(slide, color: RGBColor = C_BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill, line=None, line_w=0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w or 1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", size=14, bold=False, color=C_TEXT,
                align=PP_ALIGN.LEFT, font=FONT, wrap=True, spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_after = Pt(4)
    p.line_spacing = spacing
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb, tf


def add_rich_paragraph(tf, text, size=14, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(6)
    run = p.runs[0]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_slide_chrome(slide, section_label="", slide_num=None, source=""):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), C_ACCENT)
    add_textbox(slide, Inches(0.55), Inches(0.18), Inches(3), Inches(0.25), DATE_LABEL,
                size=9, color=C_BORDER, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.18), Inches(2.3), Inches(0.25), "CONFIDENTIAL",
                size=9, color=C_BORDER, bold=True, align=PP_ALIGN.RIGHT)
    if section_label:
        add_textbox(slide, Inches(0.55), Inches(0.55), Inches(12), Inches(0.35), section_label.upper(),
                    size=11, bold=True, color=C_CYAN)
    if source:
        add_textbox(slide, Inches(0.55), Inches(7.05), Inches(12), Inches(0.35), f"Source: {source}",
                    size=8, color=C_BORDER)
    if slide_num is not None:
        add_textbox(slide, Inches(12.4), Inches(7.05), Inches(0.6), Inches(0.25), str(slide_num),
                    size=8, color=C_BORDER, align=PP_ALIGN.RIGHT)


def add_title_block(slide, title, intro="", y=Inches(0.95)):
    add_textbox(slide, Inches(0.55), y, Inches(12.2), Inches(1.1), title,
                size=32, bold=True, color=C_WHITE, spacing=1.05)
    if intro:
        add_textbox(slide, Inches(0.55), y + Inches(1.05), Inches(11.8), Inches(0.7), intro,
                    size=14, color=C_TEXT, spacing=1.25)


def add_exit_line(slide, text):
    add_rect(slide, Inches(0.55), Inches(6.55), Inches(0.08), Inches(0.55), C_ACCENT)
    add_textbox(slide, Inches(0.75), Inches(6.55), Inches(11.8), Inches(0.55), text,
                size=13, bold=True, color=C_TEXT_BRIGHT)


def add_kpi_card(slide, left, top, width, height, value, label, sub="", accent=C_ACCENT):
    add_rounded(slide, left, top, width, height, C_SURFACE, C_BORDER)
    add_rect(slide, left, top, width, Inches(0.06), accent)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.18), width - Inches(0.4), Inches(0.7),
                value, size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_textbox(slide, left + Inches(0.2), top + height - Inches(0.85), width - Inches(0.4), Inches(0.55),
                label, size=11, color=C_TEXT, align=PP_ALIGN.LEFT)
    if sub:
        add_textbox(slide, left + Inches(0.2), top + height - Inches(0.4), width - Inches(0.4), Inches(0.3),
                    sub, size=9, color=C_BORDER)


def add_numbered_item(slide, left, top, num, title, body, width=Inches(5.5)):
    add_textbox(slide, left, top, Inches(0.45), Inches(0.35), f"{num:02d}", size=14, bold=True, color=C_GOLD)
    add_textbox(slide, left + Inches(0.45), top, width - Inches(0.45), Inches(0.35), title,
                size=13, bold=True, color=C_WHITE)
    add_textbox(slide, left + Inches(0.45), top + Inches(0.32), width - Inches(0.45), Inches(0.55), body,
                size=11, color=C_TEXT)


def add_flow_step(slide, left, top, num, title, width=Inches(1.55)):
    add_rounded(slide, left, top, width, Inches(0.95), C_SURFACE, C_BORDER)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.08), width - Inches(0.24), Inches(0.25),
                f"{num}", size=10, bold=True, color=C_GOLD)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.32), width - Inches(0.24), Inches(0.55),
                title, size=10, bold=True, color=C_WHITE)


def add_arrow_right(slide, left, top):
    add_textbox(slide, left, top, Inches(0.35), Inches(0.35), "→", size=18, bold=True, color=C_CYAN, align=PP_ALIGN.CENTER)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    n = 0

    # --- SLIDE 1: COVER ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), C_ACCENT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, C_ACCENT)
    add_textbox(slide, Inches(0.75), Inches(0.45), Inches(8), Inches(0.35),
                "FOR OWNERS:", size=12, bold=True, color=C_CYAN)
    add_textbox(slide, Inches(0.75), Inches(1.15), Inches(11.5), Inches(1.8),
                "Your Next Guest May Choose Their Hotel\nBefore They Ever Search for It",
                size=36, bold=True, color=C_WHITE, spacing=1.05)
    add_textbox(slide, Inches(0.75), Inches(3.05), Inches(10), Inches(0.45),
                "Introducing Dealality AI Demand Positioning™",
                size=20, bold=True, color=C_GOLD)
    add_textbox(slide, Inches(0.75), Inches(3.65), Inches(10.5), Inches(0.55),
                "AI is rapidly becoming a new layer between traveler intent and hotel consideration.",
                size=14, color=C_TEXT)
    add_rounded(slide, Inches(0.75), Inches(4.35), Inches(11.5), Inches(2.05), C_SURFACE, C_BORDER)
    add_textbox(slide, Inches(1.1), Inches(4.55), Inches(3.5), Inches(0.9), "56%",
                size=64, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.5), Inches(4.55), Inches(7.2), Inches(0.55),
                "of U.S. travelers used AI for planning, booking or in-destination assistance\nfor at least one trip in the past 12 months.",
                size=13, color=C_TEXT_BRIGHT)
    add_textbox(slide, Inches(4.5), Inches(5.25), Inches(7.2), Inches(0.35),
                "43% — 2H 2025  ·  33% — 1H 2025  ·  +23 pts / ~70% growth in roughly one year",
                size=11, bold=True, color=C_GREEN)
    add_textbox(slide, Inches(4.5), Inches(5.65), Inches(7.2), Inches(0.35),
                "Phocuswright: travel's \"fastest behavioral shift in a decade.\"",
                size=10, color=C_BORDER)
    add_textbox(slide, Inches(0.75), Inches(6.65), Inches(11), Inches(0.35),
                "The question is no longer whether travelers will use AI — but whether your hotels are winning when they do.",
                size=12, bold=True, color=C_TEXT_BRIGHT)
    add_textbox(slide, Inches(0.75), Inches(7.05), Inches(10), Inches(0.25),
                "Source: Phocuswright, The AI Surge: Travel's Fastest Behavioral Shift in a Decade, 2026.",
                size=8, color=C_BORDER)
    add_textbox(slide, Inches(9.5), Inches(0.45), Inches(3.2), Inches(0.35), DATE_LABEL,
                size=10, color=C_BORDER, align=PP_ALIGN.RIGHT)
    add_rounded(slide, Inches(9.2), Inches(6.55), Inches(3.5), Inches(0.45), C_SURFACE2, C_ACCENT)
    add_textbox(slide, Inches(9.35), Inches(6.62), Inches(3.2), Inches(0.3),
                "OWNER INTRODUCTION", size=9, bold=True, color=C_CYAN, align=PP_ALIGN.CENTER)

    # --- SLIDE 2: THE SHIFT ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "THE SHIFT", n,
                     "Adobe Digital Insights, June 2026; Booking.com Global AI Sentiment Report, July 2025.")
    add_title_block(slide, "Hotel Discovery Is Being Rewritten",
                    "The shift isn't theoretical. AI-driven travel discovery is already scaling at extraordinary rates.")
    add_kpi_card(slide, Inches(0.55), Inches(2.35), Inches(3.85), Inches(1.55), "+2,215%",
                 "Growth in AI-sourced traffic to U.S. travel sites\n(Oct 2024 – May 2026)", accent=C_GREEN)
    add_kpi_card(slide, Inches(4.65), Inches(2.35), Inches(3.85), Inches(1.55), "+194%",
                 "YoY growth in AI-referred travel traffic\n(May 2026 alone)", accent=C_CYAN)
    add_kpi_card(slide, Inches(8.75), Inches(2.35), Inches(3.85), Inches(1.55), "89%",
                 "Of consumers globally want to use AI\nin future travel planning", accent=C_GOLD)
    flow_y = Inches(4.15)
    add_rounded(slide, Inches(0.55), flow_y, Inches(5.5), Inches(1.15), C_SURFACE2, C_BORDER)
    add_textbox(slide, Inches(0.75), flow_y + Inches(0.12), Inches(1.2), Inches(0.25), "SEARCH",
                size=10, bold=True, color=C_BORDER)
    add_textbox(slide, Inches(0.75), flow_y + Inches(0.42), Inches(5.1), Inches(0.55),
                "Google  →  results  →  hotel", size=14, bold=True, color=C_TEXT)
    add_textbox(slide, Inches(6.35), flow_y + Inches(0.35), Inches(0.5), Inches(0.4), "→",
                size=28, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_rounded(slide, Inches(7.05), flow_y, Inches(5.55), Inches(1.15), C_SURFACE, C_ACCENT)
    add_textbox(slide, Inches(7.25), flow_y + Inches(0.12), Inches(1.5), Inches(0.25), "CONVERSATION",
                size=10, bold=True, color=C_CYAN)
    add_textbox(slide, Inches(7.25), flow_y + Inches(0.42), Inches(5.1), Inches(0.55),
                "Traveler intent  →  AI interpretation  →  recommendation  →  consideration  →  hotel/OTA",
                size=12, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.55), Inches(5.45), Inches(12), Inches(0.35),
                "Adobe: 8+ million visits to U.S. travel sites. Booking.com: 37,000+ consumers across 33 markets.",
                size=10, color=C_BORDER)
    add_exit_line(slide, "A new intermediary is forming between traveler demand and the hotels that get considered.")

    # --- SLIDE 3: WHY THIS MATTERS ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "WHY THIS MATTERS", n,
                     "Expedia Group AI Trust Gap, April 2026; Adobe Digital Insights, June 2026.")
    add_title_block(slide, "AI Is Becoming a Consideration Engine — Before It Becomes a Booking Engine",
                    "AI's influence today is disproportionately concentrated in research, discovery and consideration.")
    stats = [("53%", "comfortable allowing AI to suggest travel options", C_CYAN),
             ("48%", "say AI saves time and helps discover places they wouldn't have found", C_GOLD),
             ("68%", "still prefer to book with a trusted travel brand vs. an AI chatbot", C_TEXT)]
    for i, (val, lbl, accent) in enumerate(stats):
        add_kpi_card(slide, Inches(0.55 + i * 4.1), Inches(2.25), Inches(3.85), Inches(1.35), val, lbl, accent=accent)
    add_rounded(slide, Inches(0.55), Inches(3.85), Inches(7.2), Inches(1.05), C_SURFACE, C_ACCENT)
    add_textbox(slide, Inches(0.8), Inches(4.05), Inches(6.7), Inches(0.35),
                "AI influences the shortlist.", size=16, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(4.45), Inches(6.7), Inches(0.35),
                "The hotel or OTA still captures the transaction.", size=14, color=C_TEXT)
    add_rounded(slide, Inches(8.0), Inches(3.85), Inches(4.55), Inches(1.05), C_SURFACE2, C_BORDER)
    add_textbox(slide, Inches(8.2), Inches(4.0), Inches(4.1), Inches(0.85),
                "AI-referred travelers (May 2026):\n+70% longer visits  ·  +21% engagement  ·  −41% bounce",
                size=11, bold=True, color=C_GREEN)
    add_textbox(slide, Inches(0.55), Inches(5.15), Inches(12), Inches(0.45),
                "The emerging battleground is not simply AI booking. It is AI consideration.",
                size=18, bold=True, color=C_GOLD)
    add_exit_line(slide, "Hotel owners currently have almost no performance system for measuring it.")

    # --- SLIDE 4: BLIND SPOT ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "THE BLIND SPOT", n, "Deloitte, 2026 Travel Industry Outlook.")
    add_title_block(slide, "We Measure Almost Everything — Except This",
                    "Hotel owners have mature systems for measuring market performance, pricing, reputation, distribution and conversion.")
    cols = [
        ("MARKET PERFORMANCE", "Occupancy · ADR · RevPAR · RevPAR Index"),
        ("REVENUE", "Pace · Pickup · Segment mix · Channel mix"),
        ("DIGITAL", "Search traffic · Paid media · Website conversion"),
        ("REPUTATION", "Ratings · Reviews · Guest sentiment"),
    ]
    for i, (head, body) in enumerate(cols):
        x = Inches(0.55 + i * 3.05)
        add_rounded(slide, x, Inches(2.35), Inches(2.85), Inches(1.55), C_SURFACE, C_BORDER)
        add_textbox(slide, x + Inches(0.15), Inches(2.5), Inches(2.55), Inches(0.35), head,
                    size=9, bold=True, color=C_CYAN)
        add_textbox(slide, x + Inches(0.15), Inches(2.9), Inches(2.55), Inches(0.85), body,
                    size=11, color=C_TEXT)
    add_rounded(slide, Inches(0.55), Inches(4.15), Inches(12.05), Inches(1.55), C_SURFACE2, C_ACCENT)
    add_textbox(slide, Inches(0.85), Inches(4.3), Inches(1.5), Inches(0.35), "BUT…",
                size=20, bold=True, color=C_GOLD)
    qs = ["Who does AI recommend?", "For what demand?", "Against whom?", "And why?"]
    for i, q in enumerate(qs):
        add_textbox(slide, Inches(2.2 + i * 2.65), Inches(4.35), Inches(2.4), Inches(0.55), q,
                    size=13, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.85), Inches(5.05), Inches(11.2), Inches(0.45),
                "Deloitte: GenAI use for trip planning tripled between 2023 and 2025 — itineraries, hotels, activities.",
                size=10, color=C_TEXT)
    add_exit_line(slide, "Dealality creates the missing performance layer.")

    # --- SLIDE 5: PRODUCT INTRO ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "THE SOLUTION", n)
    add_textbox(slide, Inches(0.55), Inches(0.95), Inches(12), Inches(0.45), "Introducing",
                size=16, color=C_TEXT)
    add_textbox(slide, Inches(0.55), Inches(1.3), Inches(12), Inches(0.65),
                "AI Demand Positioning™", size=36, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.55), Inches(2.0), Inches(11.5), Inches(0.55),
                "Dealality measures how a hotel competes inside AI-driven demand — not simply whether its name appears in an AI response.",
                size=13, color=C_TEXT)
    questions = [
        ("ARE WE CONSIDERED?", "How often does the hotel enter relevant AI consideration sets?"),
        ("FOR WHAT DEMAND?", "Families? Business? Lifestyle? Groups? Experiences? Location?"),
        ("AGAINST WHOM?", "Which hotels repeatedly compete for that demand?"),
        ("WHY DO WE WIN OR LOSE?", "Product? Positioning? Reviews? Narrative? Amenities? Sources?"),
        ("WHAT SHOULD CHANGE?", "Where is the highest-value opportunity to improve?"),
    ]
    for i, (t, b) in enumerate(questions):
        row, col = divmod(i, 2)
        add_numbered_item(slide, Inches(0.55 + col * 6.2), Inches(2.65 + row * 0.95), i + 1, t, b, Inches(5.8))
    add_rounded(slide, Inches(0.55), Inches(5.35), Inches(7.5), Inches(0.95), C_SURFACE, C_BORDER)
    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(3.2), Inches(0.3), "FROM AI VISIBILITY",
                size=10, bold=True, color=C_BORDER)
    add_textbox(slide, Inches(0.8), Inches(5.78), Inches(3.2), Inches(0.35), "\"Did AI mention my hotel?\"",
                size=12, color=C_TEXT)
    add_textbox(slide, Inches(4.0), Inches(5.65), Inches(0.4), Inches(0.4), "↓", size=22, bold=True, color=C_ACCENT)
    add_textbox(slide, Inches(4.45), Inches(5.5), Inches(3.5), Inches(0.3), "TO AI DEMAND INTELLIGENCE",
                size=10, bold=True, color=C_CYAN)
    add_textbox(slide, Inches(4.45), Inches(5.78), Inches(3.5), Inches(0.35),
                "\"Where are we winning and losing demand — and why?\"", size=12, bold=True, color=C_WHITE)
    add_rounded(slide, Inches(8.3), Inches(5.35), Inches(4.3), Inches(0.95), C_SURFACE2, C_GOLD)
    add_textbox(slide, Inches(8.5), Inches(5.5), Inches(3.9), Inches(0.65),
                "Nearly 4 in 10 U.S. travelers used GenAI while researching trips in 2025 — Phocuswright",
                size=10, color=C_TEXT_BRIGHT)
    add_exit_line(slide, "To do that, Dealality observes demand at a scale that a human analyst cannot.")

    # --- SLIDE 6: HOW IT WORKS ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "HOW IT WORKS", n)
    add_title_block(slide, "An Always-On AI Demand Intelligence Engine",
                    "Dealality converts thousands of individual AI interactions into structured, repeatable hotel intelligence.")
    steps = [
        ("DEFINE\nDEMAND", "Hundreds or thousands of relevant traveler intentions"),
        ("OBSERVE", "Run structured demand scenarios across AI environments"),
        ("CAPTURE", "Recommendations · Competitors · Position · Narratives · Sources"),
        ("NORMALIZE", "Turn individual AI responses into comparable observations"),
        ("BENCHMARK", "Hotel · Comp set · Market · Portfolio"),
        ("DIAGNOSE", "Identify why demand is being won or lost"),
        ("RETEST", "Measure change over time"),
    ]
    x = Inches(0.45)
    for i, (title, _) in enumerate(steps):
        add_flow_step(slide, x, Inches(2.35), i + 1, title)
        if i < len(steps) - 1:
            add_arrow_right(slide, x + Inches(1.58), Inches(2.65))
        x += Inches(1.78)
    add_rounded(slide, Inches(0.55), Inches(3.55), Inches(12.05), Inches(2.35), C_SURFACE, C_BORDER)
    details_y = Inches(3.75)
    for i, (_, body) in enumerate(steps):
        col = i % 4
        row = i // 4
        add_textbox(slide, Inches(0.8 + col * 3.0), details_y + Inches(row * 0.95), Inches(2.75), Inches(0.8),
                    body, size=10, color=C_TEXT)
    add_textbox(slide, Inches(0.55), Inches(6.15), Inches(12), Inches(0.35),
                "One prompt is anecdotal. Repeated observation becomes intelligence.",
                size=16, bold=True, color=C_GOLD)
    add_exit_line(slide, "A measurable AI consideration position — something that does not exist in traditional hotel reporting.")

    # --- SLIDE 7: AI CONSIDERATION INDEX ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "THE KPI", n)
    add_title_block(slide, "A New Performance Metric: AI Consideration Index™",
                    "We want AI demand performance to be as intuitive to an owner as RevPAR Index.")
    add_rounded(slide, Inches(3.6), Inches(2.2), Inches(6.0), Inches(2.0), C_SURFACE, C_ACCENT)
    add_textbox(slide, Inches(3.85), Inches(2.35), Inches(5.5), Inches(0.35),
                "AI CONSIDERATION INDEX", size=12, bold=True, color=C_CYAN, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.85), Inches(2.75), Inches(5.5), Inches(0.9), "112",
                size=72, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.85), Inches(3.65), Inches(5.5), Inches(0.35),
                "12% ABOVE EXPECTED COMPETITIVE SHARE", size=13, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.85), Inches(4.0), Inches(5.5), Inches(0.35),
                "100 = Expected share  ·  >100 = Outperforming  ·  <100 = Underperforming",
                size=9, color=C_BORDER, align=PP_ALIGN.CENTER)
    kpis = [("67%", "Consideration Rate"), ("54%", "Competitive Win Rate"), ("41%", "Top Recommendation Rate"),
            ("+8 pts", "90-Day Position Change"), ("72%", "Relevant Demand Coverage")]
    for i, (v, l) in enumerate(kpis):
        add_kpi_card(slide, Inches(0.55 + i * 2.45), Inches(4.45), Inches(2.25), Inches(1.15), v, l, accent=C_SURFACE2)
    add_textbox(slide, Inches(0.55), Inches(5.85), Inches(12), Inches(0.3),
                "Figures shown are illustrative examples of Dealality product output, not third-party market statistics.",
                size=9, color=C_BORDER)
    add_exit_line(slide, "An overall score can hide the most important insight: which demand the hotel is actually winning.")

    # --- SLIDE 8: PRODUCT DEMO ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "LIVE PRODUCT", n)
    add_title_block(slide, "AI Demand Positioning — Product View",
                    "Illustrative dashboard showing how owners see consideration performance in one place.")
    add_rounded(slide, Inches(0.55), Inches(2.05), Inches(12.05), Inches(4.55), C_SURFACE2, C_ACCENT)
    add_rect(slide, Inches(0.55), Inches(2.05), Inches(12.05), Inches(0.45), C_SURFACE)
    add_textbox(slide, Inches(0.8), Inches(2.12), Inches(4), Inches(0.3), "DEALALITY  ·  AI DEMAND POSITIONING",
                size=10, bold=True, color=C_CYAN)
    add_textbox(slide, Inches(10.5), Inches(2.12), Inches(2), Inches(0.3), "Waterstone Boca Raton",
                size=9, color=C_TEXT, align=PP_ALIGN.RIGHT)
    demo_kpis = [("112", "AI Index", C_ACCENT), ("67%", "Consideration", C_CYAN), ("54%", "Win Rate", C_GREEN),
                 ("31%", "Family Gap", C_RED), ("+8 pts", "90-Day Δ", C_GOLD)]
    for i, (v, l, accent) in enumerate(demo_kpis):
        add_kpi_card(slide, Inches(0.8 + i * 2.35), Inches(2.65), Inches(2.15), Inches(1.05), v, l, accent=accent)
    add_rounded(slide, Inches(0.8), Inches(3.95), Inches(5.6), Inches(2.35), C_BG, C_BORDER)
    add_textbox(slide, Inches(1.0), Inches(4.05), Inches(5.2), Inches(0.3), "DEMAND POSITION MAP",
                size=10, bold=True, color=C_CYAN)
    rows = [("Business Travel", "78%", "+18 pts", "↑"), ("Waterfront", "74%", "+14 pts", "↑"),
            ("Couples", "62%", "+3 pts", "→"), ("Families", "31%", "−24 pts", "↓")]
    y = Inches(4.4)
    for terr, cons, gap, trend in rows:
        add_textbox(slide, Inches(1.0), y, Inches(2.2), Inches(0.28), terr, size=10, color=C_TEXT)
        add_textbox(slide, Inches(3.2), y, Inches(0.8), Inches(0.28), cons, size=10, bold=True, color=C_WHITE)
        add_textbox(slide, Inches(4.0), y, Inches(1.0), Inches(0.28), gap, size=10,
                    color=C_GREEN if "+" in gap else C_RED)
        add_textbox(slide, Inches(5.1), y, Inches(0.4), Inches(0.28), trend, size=10, bold=True, color=C_CYAN)
        y += Inches(0.32)
    add_rounded(slide, Inches(6.65), Inches(3.95), Inches(5.75), Inches(2.35), C_BG, C_BORDER)
    add_textbox(slide, Inches(6.85), Inches(4.05), Inches(5.3), Inches(0.3), "PRIORITY OPPORTUNITY — FAMILY DEMAND",
                size=10, bold=True, color=C_GOLD)
    add_textbox(slide, Inches(6.85), Inches(4.45), Inches(2.5), Inches(0.55),
                "Consideration\n31%", size=18, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(9.5), Inches(4.45), Inches(2.5), Inches(0.55),
                "Competitive Gap\n−24 pts", size=18, bold=True, color=C_RED)
    add_textbox(slide, Inches(6.85), Inches(5.25), Inches(5.3), Inches(0.85),
                "Limited AI evidence: room configurations · family activities · pool/recreation · local family experiences",
                size=10, color=C_TEXT)
    add_exit_line(slide, "From headline index to demand territories, competitors, and actionable opportunities — in one platform.")

    # --- SLIDE 9: OBSERVED DEMAND ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "OBSERVED DEMAND", n, "Booking.com Global AI Sentiment Report, July 2025.")
    add_title_block(slide, "Travelers Don't Search in Hotel Segments. They Express Needs.",
                    "Conversational AI allows travelers to communicate far more context than traditional keyword search.")
    prompts = [
        "\"Upscale hotel in Tampa near the waterfront for a family with two children.\"",
        "\"Cool hotel in Waikiki that feels local rather than corporate.\"",
        "\"Best San Francisco hotel for a three-day business trip near Union Square.\"",
        "\"Boutique hotel in Newport for a couples weekend with good restaurants nearby.\"",
    ]
    for i, p in enumerate(prompts):
        add_rounded(slide, Inches(0.55), Inches(2.15 + i * 0.52), Inches(6.0), Inches(0.42), C_SURFACE, C_BORDER)
        add_textbox(slide, Inches(0.75), Inches(2.22 + i * 0.52), Inches(5.6), Inches(0.3), p, size=10, color=C_TEXT_BRIGHT)
    territories = [
        ("WHO", "Families · Couples · Business · Groups · Luxury travelers"),
        ("WHY", "Business · Weekend escape · Event · Family vacation · Experience"),
        ("NEED", "Walkability · Dining · Pool · Rooms · Meetings · Design · Location"),
    ]
    for i, (head, body) in enumerate(territories):
        add_rounded(slide, Inches(6.85), Inches(2.15 + i * 0.95), Inches(5.75), Inches(0.8), C_SURFACE2,
                    C_ACCENT if i == 0 else C_BORDER)
        add_textbox(slide, Inches(7.05), Inches(2.25 + i * 0.95), Inches(1.0), Inches(0.3), head,
                    size=11, bold=True, color=C_GOLD)
        add_textbox(slide, Inches(7.85), Inches(2.25 + i * 0.95), Inches(4.5), Inches(0.55), body,
                    size=11, color=C_TEXT)
    booking_stats = [("38%", "research destinations / timing"), ("37%", "discover local experiences"),
                     ("36%", "find restaurant recommendations")]
    for i, (v, l) in enumerate(booking_stats):
        add_kpi_card(slide, Inches(0.55 + i * 4.1), Inches(4.35), Inches(3.85), Inches(1.05), v, l, accent=C_CYAN)
    add_textbox(slide, Inches(0.55), Inches(5.65), Inches(12), Inches(0.4),
                "AI gives us a new way to observe not only search volume — but demand intent.",
                size=16, bold=True, color=C_GOLD)
    add_exit_line(slide, "A hotel can be strong overall while being nearly invisible for a strategically important guest need.")

    # --- SLIDE 10: DEMAND POSITION MAP ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "DEMAND POSITION MAP", n)
    add_title_block(slide, "See Where Each Asset Wins — and Where It Disappears",
                    "Dealality breaks overall AI performance into the demand territories that matter economically to the hotel.")
    headers = ["Demand Territory", "Consideration", "vs. Competitors", "Trend"]
    col_x = [Inches(0.75), Inches(4.0), Inches(6.2), Inches(8.6)]
    add_rect(slide, Inches(0.55), Inches(2.15), Inches(8.7), Inches(0.4), C_SURFACE2)
    for h, x in zip(headers, col_x):
        add_textbox(slide, x, Inches(2.2), Inches(2.2), Inches(0.3), h, size=9, bold=True, color=C_CYAN)
    table_rows = [
        ("Business Travel", "78%", "+18 pts", "↑"), ("Waterfront", "74%", "+14 pts", "↑"),
        ("Couples", "62%", "+3 pts", "→"), ("Groups", "58%", "−2 pts", "→"),
        ("Families", "31%", "−24 pts", "↓"), ("Local Experience", "27%", "−29 pts", "↓"),
    ]
    for ri, row in enumerate(table_rows):
        y = Inches(2.6 + ri * 0.38)
        bg = C_SURFACE if ri % 2 == 0 else C_BG
        add_rect(slide, Inches(0.55), y - Inches(0.04), Inches(8.7), Inches(0.36), bg)
        colors = [C_TEXT, C_WHITE, C_GREEN if "+" in row[2] else C_RED, C_CYAN]
        for ci, (cell, x) in enumerate(zip(row, col_x)):
            add_textbox(slide, x, y, Inches(2.2), Inches(0.3), cell, size=10,
                        bold=ci == 0, color=colors[ci])
    add_rounded(slide, Inches(9.5), Inches(2.15), Inches(3.1), Inches(3.55), C_SURFACE, C_RED)
    add_textbox(slide, Inches(9.7), Inches(2.3), Inches(2.7), Inches(0.3), "FAMILY DEMAND — DIAGNOSIS",
                size=10, bold=True, color=C_GOLD)
    add_textbox(slide, Inches(9.7), Inches(2.65), Inches(2.7), Inches(0.55),
                "Consideration: 31%\nCompetitive gap: −24 pts", size=12, bold=True, color=C_WHITE)
    reasons = [
        "Competitor A — larger rooms narrative",
        "Competitor B — strong pool/family story",
        "Competitor C — walkable family activities",
    ]
    for i, r in enumerate(reasons):
        add_textbox(slide, Inches(9.7), Inches(3.35 + i * 0.42), Inches(2.7), Inches(0.35), f"• {r}",
                    size=9, color=C_TEXT)
    add_textbox(slide, Inches(0.55), Inches(5.15), Inches(12), Inches(0.35),
                "This is not simply an SEO problem. It is a demand positioning problem. (Product figures illustrative.)",
                size=13, bold=True, color=C_GOLD)
    add_exit_line(slide, "The hotels taking that demand may not be the hotels management currently considers competitors.")

    # --- SLIDE 11: DYNAMIC COMP SET ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "DYNAMIC COMP SET", n)
    add_title_block(slide, "AI Does Not Know Your STR Comp Set",
                    "When a traveler asks AI for a recommendation, the model constructs a competitive set dynamically around the traveler's need.")
    add_rounded(slide, Inches(0.55), Inches(2.25), Inches(5.5), Inches(2.2), C_SURFACE2, C_BORDER)
    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(5), Inches(0.3), "TRADITIONAL COMP SET",
                size=11, bold=True, color=C_BORDER)
    add_textbox(slide, Inches(0.8), Inches(2.85), Inches(5), Inches(0.35), "Hotel", size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(3.35), Inches(5), Inches(0.3), "versus", size=12, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(3.65), Inches(5), Inches(0.55),
                "5–8 predetermined competitors", size=14, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_rounded(slide, Inches(6.55), Inches(2.25), Inches(6.05), Inches(2.2), C_SURFACE, C_ACCENT)
    add_textbox(slide, Inches(6.8), Inches(2.4), Inches(5.5), Inches(0.3), "DEALALITY OBSERVED COMP SET™",
                size=11, bold=True, color=C_CYAN)
    add_textbox(slide, Inches(6.8), Inches(2.85), Inches(5.5), Inches(0.35), "Hotel", size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(3.35), Inches(5.5), Inches(0.3), "versus", size=12, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.8), Inches(3.65), Inches(5.5), Inches(0.55),
                "Every property repeatedly entering the same AI demand conversations",
                size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tags = [("EXPECTED", "Hotels already in management's competitive set", C_TEXT),
            ("EMERGING", "Properties appearing with increasing frequency", C_CYAN),
            ("UNEXPECTED", "Hotels outside the traditional comp set taking consideration", C_GOLD),
            ("DEMAND-SPECIFIC", "Different competitors for families vs. groups vs. business", C_GREEN)]
    for i, (t, b, c) in enumerate(tags):
        row, col = divmod(i, 2)
        add_rounded(slide, Inches(0.55 + col * 6.2), Inches(4.65 + row * 0.85), Inches(5.95), Inches(0.72), C_SURFACE2, c)
        add_textbox(slide, Inches(0.75 + col * 6.2), Inches(4.75 + row * 0.85), Inches(1.4), Inches(0.25), t,
                    size=9, bold=True, color=c)
        add_textbox(slide, Inches(2.1 + col * 6.2), Inches(4.75 + row * 0.85), Inches(4.2), Inches(0.45), b,
                    size=10, color=C_TEXT)
    add_textbox(slide, Inches(0.55), Inches(6.15), Inches(12), Inches(0.35),
                "The competitive set becomes an observed outcome — not only a management assumption.",
                size=15, bold=True, color=C_GOLD)
    add_exit_line(slide, "Knowing the competitor is useful. Understanding why the competitor wins is significantly more valuable.")

    # --- SLIDE 12: WIN / LOSS ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "WIN / LOSS", n, "Phocuswright, The AI Surge, 2026.")
    add_title_block(slide, "Move From \"Who Won?\" to \"Why Did They Win?\"",
                    "Dealality analyzes recurring recommendation patterns to identify the attributes, narratives and evidence driving competitive outcomes.")
    add_rounded(slide, Inches(0.55), Inches(2.25), Inches(5.8), Inches(2.85), C_SURFACE, C_GREEN)
    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(5.3), Inches(0.3), "YOUR HOTEL — WINS",
                size=11, bold=True, color=C_GREEN)
    for i, item in enumerate(["+ Location", "+ Meeting capability", "+ Brand awareness", "+ Waterfront"]):
        add_textbox(slide, Inches(0.8), Inches(2.75 + i * 0.35), Inches(5.3), Inches(0.3), item, size=11, color=C_TEXT_BRIGHT)
    add_textbox(slide, Inches(0.8), Inches(4.15), Inches(5.3), Inches(0.3), "YOUR HOTEL — LOSES",
                size=11, bold=True, color=C_RED)
    for i, item in enumerate(["− Family suitability", "− Local character", "− F&B narrative", "− Distinctive experience"]):
        add_textbox(slide, Inches(0.8), Inches(4.5 + i * 0.35), Inches(5.3), Inches(0.3), item, size=11, color=C_TEXT)
    add_rounded(slide, Inches(6.65), Inches(2.25), Inches(5.85), Inches(2.85), C_SURFACE2, C_BORDER)
    add_textbox(slide, Inches(6.9), Inches(2.4), Inches(5.3), Inches(0.35),
                "COMPETITOR A — Wins 68% of head-to-head family scenarios", size=12, bold=True, color=C_WHITE)
    win_reasons = [("42%", "room/family configuration narrative"), ("31%", "pool/recreation"),
                   ("28%", "nearby family experiences"), ("24%", "review language")]
    for i, (pct, reason) in enumerate(win_reasons):
        add_textbox(slide, Inches(6.9), Inches(2.95 + i * 0.45), Inches(1.0), Inches(0.35), pct,
                    size=14, bold=True, color=C_GOLD)
        add_textbox(slide, Inches(7.85), Inches(2.95 + i * 0.45), Inches(4.3), Inches(0.35), reason,
                    size=11, color=C_TEXT)
    add_rounded(slide, Inches(0.55), Inches(5.35), Inches(12.05), Inches(0.65), C_SURFACE2, C_GOLD)
    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.35),
                "Phocuswright: AI-using travelers have higher household incomes, take more trips, and spend more annually on travel.",
                size=11, color=C_TEXT_BRIGHT)
    add_exit_line(slide, "The next question is what is teaching AI to make those distinctions.")

    # --- SLIDE 13: NARRATIVE + SOURCE ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "NARRATIVE + SOURCE INTELLIGENCE", n, "Adobe Digital Insights, June 2026.")
    add_title_block(slide, "What Does AI Believe About Your Hotel — and Why?",
                    "AI does not rely exclusively on the hotel's own marketing message. Its recommendations synthesize information across a broad digital ecosystem.")
    narrative_cols = [
        ("WHAT YOU SAY", "Hotel website\nBrand content\nProperty descriptions", C_CYAN),
        ("WHAT THE MARKET SAYS", "Guest reviews\nTravel publications\nDestination content\nSocial conversation\nThird-party travel sites", C_GOLD),
        ("WHAT AI BELIEVES", "Who the hotel is for\nWhat makes it different\nWhere it wins\nWhen it should be recommended", C_GREEN),
    ]
    for i, (head, body, accent) in enumerate(narrative_cols):
        x = Inches(0.55 + i * 4.15)
        add_rounded(slide, x, Inches(2.25), Inches(3.85), Inches(2.35), C_SURFACE, accent)
        add_textbox(slide, x + Inches(0.2), Inches(2.4), Inches(3.45), Inches(0.35), head,
                    size=11, bold=True, color=accent)
        add_textbox(slide, x + Inches(0.2), Inches(2.85), Inches(3.45), Inches(1.55), body, size=11, color=C_TEXT)
        if i < 2:
            add_textbox(slide, x + Inches(3.95), Inches(3.15), Inches(0.35), Inches(0.35), "→",
                        size=20, bold=True, color=C_ACCENT)
    add_textbox(slide, Inches(0.55), Inches(4.85), Inches(12), Inches(0.45),
                "Adobe 2026: More than one-third of key hotel pages remain unreadable by AI systems.",
                size=11, color=C_BORDER)
    analysis = ["Narrative Strength", "Narrative Gaps", "Narrative Conflicts", "Source Influence",
                "Review Signals", "Social Sentiment"]
    for i, a in enumerate(analysis):
        add_rounded(slide, Inches(0.55 + i * 2.05), Inches(5.4), Inches(1.85), Inches(0.45), C_SURFACE2, C_BORDER)
        add_textbox(slide, Inches(0.6 + i * 2.05), Inches(5.48), Inches(1.75), Inches(0.3), a,
                    size=8, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.55), Inches(6.05), Inches(12), Inches(0.35),
                "Hotel positioning is no longer defined only by what you publish.",
                size=15, bold=True, color=C_GOLD)
    add_exit_line(slide, "Dealality turns the gap between intended positioning and observed positioning into an opportunity map.")

    # --- SLIDE 14: OPPORTUNITY ENGINE ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "OPPORTUNITY ENGINE", n)
    add_title_block(slide, "Don't Give Me Another Dashboard. Tell Me What to Do.",
                    "The value of intelligence is not another score — it is identifying the actions most likely to change competitive position.")
    add_rounded(slide, Inches(0.55), Inches(2.2), Inches(7.5), Inches(3.85), C_SURFACE, C_ACCENT)
    add_textbox(slide, Inches(0.8), Inches(2.35), Inches(2), Inches(0.3), "#1 OPPORTUNITY",
                size=10, bold=True, color=C_GOLD)
    add_textbox(slide, Inches(0.8), Inches(2.65), Inches(3), Inches(0.4), "FAMILY DEMAND",
                size=20, bold=True, color=C_WHITE)
    add_kpi_card(slide, Inches(0.8), Inches(3.15), Inches(2.2), Inches(1.0), "31%", "Current Consideration", accent=C_RED)
    add_kpi_card(slide, Inches(3.2), Inches(3.15), Inches(2.2), Inches(1.0), "−24 pts", "Competitive Gap", accent=C_RED)
    add_rounded(slide, Inches(5.6), Inches(3.15), Inches(2.0), Inches(1.0), C_GREEN)
    add_textbox(slide, Inches(5.75), Inches(3.35), Inches(1.7), Inches(0.55), "HIGH\nOPPORTUNITY",
                size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(4.35), Inches(3.2), Inches(0.25), "WHY", size=10, bold=True, color=C_CYAN)
    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(3.2), Inches(1.2),
                "Limited AI evidence:\n• Room configurations\n• Family-friendly activities\n• Pool / recreation\n• Family local experiences",
                size=10, color=C_TEXT)
    add_textbox(slide, Inches(4.2), Inches(4.35), Inches(3.6), Inches(0.25), "WHAT TO DO", size=10, bold=True, color=C_CYAN)
    actions = ["1. PRODUCT — Surface room configurations", "2. CONTENT — Build family narratives",
               "3. DISTRIBUTION — Correct attribute gaps", "4. REPUTATION — Strengthen guest narrative",
               "5. AI PRESENCE — Support authoritative sources"]
    for i, a in enumerate(actions):
        add_textbox(slide, Inches(4.2), Inches(4.6 + i * 0.32), Inches(3.6), Inches(0.28), a, size=9, color=C_TEXT)
    add_rounded(slide, Inches(8.3), Inches(2.2), Inches(4.3), Inches(3.85), C_SURFACE2, C_BORDER)
    add_textbox(slide, Inches(8.5), Inches(2.35), Inches(3.9), Inches(0.3), "THEN RETEST",
                size=11, bold=True, color=C_GOLD)
    add_textbox(slide, Inches(8.5), Inches(2.85), Inches(3.9), Inches(0.55),
                "Baseline  →  intervention  →  observed movement", size=14, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(8.5), Inches(3.65), Inches(3.9), Inches(1.5),
                "For an owner, this stops being a marketing exercise and becomes an asset-positioning tool.",
                size=12, color=C_TEXT)
    add_exit_line(slide, "Actionable intelligence — not another score without a path forward.")

    # --- SLIDE 15: NEWBOND + DOVETAIL ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "WHY NEWBOND + DOVETAIL", n)
    add_title_block(slide, "Built for the Questions Hotel Owners Actually Ask",
                    "Particularly powerful across portfolios where repositioning, differentiated experiences and active asset management drive value.")
    add_rounded(slide, Inches(0.55), Inches(2.25), Inches(5.95), Inches(3.35), C_SURFACE, C_BORDER)
    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(0.35), "NEWBOND",
                size=16, bold=True, color=C_CYAN)
    add_kpi_card(slide, Inches(0.8), Inches(2.85), Inches(1.7), Inches(1.05), "15+", "Hotels", accent=C_ACCENT)
    add_kpi_card(slide, Inches(2.65), Inches(2.85), Inches(1.7), Inches(1.05), "6,000+", "Keys", accent=C_ACCENT)
    add_kpi_card(slide, Inches(4.5), Inches(2.85), Inches(1.7), Inches(1.05), "$1.0B", "Budgeted annual revenue", accent=C_GOLD)
    add_textbox(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(1.25),
                "Opportunistic acquisition and repositioning of upper-upscale and luxury full-service hotels — making observed demand positioning particularly relevant.",
                size=11, color=C_TEXT)
    add_rounded(slide, Inches(6.75), Inches(2.25), Inches(5.85), Inches(3.35), C_SURFACE, C_BORDER)
    add_textbox(slide, Inches(7.0), Inches(2.4), Inches(5.3), Inches(0.35), "DOVETAIL + CO",
                size=16, bold=True, color=C_CYAN)
    add_textbox(slide, Inches(7.0), Inches(2.9), Inches(5.3), Inches(0.35),
                "Owner and creator of experiential real estate", size=12, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(7.0), Inches(3.35), Inches(5.3), Inches(0.35),
                "New York · Newport · Waikiki · Bishop · Bermuda", size=12, color=C_GOLD)
    add_textbox(slide, Inches(7.0), Inches(3.85), Inches(5.3), Inches(1.25),
                "Explicit emphasis on hotels whose story and design create differentiated experiences.",
                size=11, color=C_TEXT)
    add_rounded(slide, Inches(0.55), Inches(5.85), Inches(12.05), Inches(0.65), C_SURFACE2, C_ACCENT)
    add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.35),
                "If differentiation creates value, you should be able to measure whether the market — and AI — actually recognizes it.",
                size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_exit_line(slide, "Significantly more valuable when monitored across an entire portfolio.")

    # --- SLIDE 16: PORTFOLIO ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "PORTFOLIO INTELLIGENCE", n)
    add_title_block(slide, "One Hotel Gives You Insight. A Portfolio Gives You Intelligence.",
                    "AI Demand Positioning creates a common framework for comparing positioning strength and opportunity across assets.")
    port_headers = ["Asset", "AI Consideration Index", "90-Day", "Largest Gap"]
    port_col_x = [Inches(0.75), Inches(3.5), Inches(6.5), Inches(8.5)]
    add_rect(slide, Inches(0.55), Inches(2.15), Inches(9.2), Inches(0.38), C_SURFACE2)
    for h, x in zip(port_headers, port_col_x):
        add_textbox(slide, x, Inches(2.2), Inches(2.5), Inches(0.28), h, size=9, bold=True, color=C_CYAN)
    portfolio = [("Hotel A", "118", "+9", "Extend leadership"), ("Hotel B", "106", "+2", "Couples"),
                 ("Hotel C", "97", "+7", "Business"), ("Hotel D", "84", "−6", "Families"),
                 ("Hotel E", "72", "−4", "Local experience")]
    for ri, row in enumerate(portfolio):
        y = Inches(2.58 + ri * 0.36)
        bg = C_SURFACE if ri % 2 == 0 else C_BG
        add_rect(slide, Inches(0.55), y - Inches(0.03), Inches(9.2), Inches(0.34), bg)
        for ci, (cell, x) in enumerate(zip(row, port_col_x)):
            color = C_GREEN if cell.startswith("+") else C_RED if cell.startswith("−") else C_WHITE if ci == 1 else C_TEXT
            add_textbox(slide, x, y, Inches(2.5), Inches(0.28), cell, size=10, bold=ci <= 1, color=color)
    use_cases = [
        ("ASSET MANAGEMENT", "Where are we losing consideration?"),
        ("REPOSITIONING", "Is market perception actually changing?"),
        ("CAPEX", "Which product gaps align with identifiable demand?"),
        ("ACQUISITIONS", "What demand position are we acquiring?"),
        ("BRANDING", "Does intended positioning match observed positioning?"),
        ("MARKETING", "Which narratives need strengthening?"),
        ("PORTFOLIO STRATEGY", "Where is the largest positioning upside?"),
    ]
    add_textbox(slide, Inches(10.0), Inches(2.15), Inches(2.8), Inches(0.3), "OWNER USE CASES",
                size=9, bold=True, color=C_GOLD)
    for i, (t, b) in enumerate(use_cases):
        add_textbox(slide, Inches(10.0), Inches(2.5 + i * 0.48), Inches(2.8), Inches(0.42),
                    f"{t}\n{b}", size=8, color=C_TEXT)
    add_textbox(slide, Inches(0.55), Inches(4.55), Inches(9), Inches(0.25),
                "Illustrative Dealality output.", size=9, color=C_BORDER)
    add_exit_line(slide, "Establish an intelligence baseline now — while AI travel behavior is still being reshaped.")

    # --- SLIDE 17: URGENCY ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "WHY NOW", n,
                     "Phocuswright 2026; Adobe Digital Insights June 2026; Deloitte 2026 Travel Industry Outlook.")
    add_title_block(slide, "The Window to Establish Position Is Now",
                    "Three things are happening simultaneously: traveler adoption is accelerating, AI referrals are accelerating, and AI systems are becoming more capable.")
    urgency = [
        ("33% → 56%", "U.S. travelers using AI around travel\n1H25 → 1H26", C_GREEN),
        ("+2,215%", "AI-sourced U.S. travel traffic\nOct 2024 → May 2026", C_CYAN),
        ("3X", "Increase in GenAI trip-planning adoption\n2023 → 2025", C_GOLD),
    ]
    for i, (v, l, accent) in enumerate(urgency):
        add_kpi_card(slide, Inches(0.55 + i * 4.1), Inches(2.35), Inches(3.85), Inches(1.65), v, l, accent=accent)
    add_rounded(slide, Inches(0.55), Inches(4.25), Inches(12.05), Inches(0.75), C_SURFACE2, C_BORDER)
    add_textbox(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.45),
                "And the technology is still early. Deloitte: adoption is occurring before fully integrated discovery and booking experiences are widely available.",
                size=12, color=C_TEXT)
    add_rounded(slide, Inches(0.55), Inches(5.25), Inches(12.05), Inches(1.05), C_SURFACE, C_ACCENT)
    add_textbox(slide, Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.55),
                "The best time to understand your AI demand position is before it becomes a board-level KPI.",
                size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_exit_line(slide, "We believe owners should start building that baseline now.")

    # --- SLIDE 18: PILOT ---
    n += 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_slide_chrome(slide, "PILOT", n)
    add_title_block(slide, "Start With the Assets Where Positioning Matters Most",
                    "Rather than a broad commitment, Dealality can demonstrate the intelligence against a small group of real hotels.")
    add_rounded(slide, Inches(0.55), Inches(2.15), Inches(3.2), Inches(0.55), C_GOLD)
    add_textbox(slide, Inches(0.7), Inches(2.28), Inches(2.9), Inches(0.3), "PROPOSED 3–5 HOTEL PILOT",
                size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    phases = [
        ("PHASE 1", "Establish Demand Universe", "What demand should each asset compete for?"),
        ("PHASE 2", "Measure Baseline", "Where does each asset currently stand?"),
        ("PHASE 3", "Discover Actual Competition", "Who repeatedly competes for that demand?"),
        ("PHASE 4", "Diagnose Win / Loss", "Why do competitors win?"),
        ("PHASE 5", "Identify Opportunities", "What can ownership, operations or marketing change?"),
        ("PHASE 6", "Retest", "Did AI demand position move?"),
    ]
    for i, (phase, title, body) in enumerate(phases):
        row, col = divmod(i, 2)
        x = Inches(0.55 + col * 6.2)
        y = Inches(2.85 + row * 1.05)
        add_rounded(slide, x, y, Inches(5.95), Inches(0.92), C_SURFACE, C_BORDER)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(1.2), Inches(0.25), phase,
                    size=9, bold=True, color=C_CYAN)
        add_textbox(slide, x + Inches(1.35), y + Inches(0.1), Inches(4.4), Inches(0.28), title,
                    size=12, bold=True, color=C_WHITE)
        add_textbox(slide, x + Inches(1.35), y + Inches(0.42), Inches(4.4), Inches(0.4), body, size=10, color=C_TEXT)
    add_textbox(slide, Inches(0.55), Inches(6.05), Inches(12), Inches(0.25), "DELIVERABLES", size=10, bold=True, color=C_GOLD)
    deliverables = ["AI Consideration Index", "Demand Position Map", "Observed Competitive Set", "Competitive Win/Loss",
                    "Narrative Intelligence", "Source Intelligence", "Social Sentiment", "Priority Opportunity Map",
                    "Portfolio Comparison"]
    for i, d in enumerate(deliverables):
        add_rounded(slide, Inches(0.55 + (i % 5) * 2.45), Inches(6.35 + (i // 5) * 0.38), Inches(2.3), Inches(0.3),
                    C_SURFACE2, C_BORDER)
        add_textbox(slide, Inches(0.6 + (i % 5) * 2.45), Inches(6.4 + (i // 5) * 0.38), Inches(2.2), Inches(0.22), d,
                    size=7, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_exit_line(slide, "Within one platform, ownership gets a completely new view of how its hotels compete for demand.")

    # --- SLIDE 19: CLOSE ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), C_ACCENT)
    add_textbox(slide, Inches(0.75), Inches(0.55), Inches(11.5), Inches(1.2),
                "Demand Is Being Shaped\nBefore the Booking Begins",
                size=34, bold=True, color=C_WHITE, spacing=1.05)
    add_textbox(slide, Inches(0.75), Inches(1.85), Inches(11), Inches(0.55),
                "The hotel industry has spent decades building systems to understand what happens after demand appears. AI gives us the opportunity to understand part of what happens before.",
                size=13, color=C_TEXT)
    close_stats = [("56%", "already using AI around travel"), ("+2,215%", "growth in AI-sourced travel traffic"),
                   ("89%", "want to use AI in future travel planning")]
    for i, (v, l) in enumerate(close_stats):
        add_kpi_card(slide, Inches(0.75 + i * 3.9), Inches(2.65), Inches(3.55), Inches(1.15), v, l, accent=C_ACCENT)
    add_rounded(slide, Inches(0.75), Inches(4.15), Inches(11.5), Inches(1.55), C_SURFACE, C_ACCENT)
    add_textbox(slide, Inches(1.0), Inches(4.45), Inches(11), Inches(0.55),
                "WHO IS AI PUTTING IN THE CONSIDERATION SET?", size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(5.15), Inches(11), Inches(0.35),
                "Dealality AI Demand Positioning™", size=16, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(5.5), Inches(11), Inches(0.3),
                "Measure it. Understand it. Improve it.", size=13, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.75), Inches(6.0), Inches(11.5), Inches(0.75),
                "You already know your RevPAR Index. You know your guest satisfaction scores. You know your channel mix. What you cannot see today is whether AI considers your hotel when a traveler describes the demand your asset was built to capture.",
                size=11, color=C_TEXT_BRIGHT)
    add_textbox(slide, Inches(0.75), Inches(6.85), Inches(5), Inches(0.25), "www.Dealality.com",
                size=10, bold=True, color=C_CYAN)
    add_textbox(slide, Inches(5.5), Inches(6.85), Inches(4), Inches(0.25), "joan@aohospitalityadvisors.com",
                size=10, color=C_TEXT)
    add_textbox(slide, Inches(9.5), Inches(6.85), Inches(2.5), Inches(0.25), "+34 674 993 637",
                size=10, color=C_TEXT, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(0.75), Inches(7.1), Inches(11), Inches(0.25),
                "CONFIDENTIAL INTRODUCTIONS AVAILABLE BY REQUEST.", size=8, color=C_BORDER)

    return prs


def resolve_out_dir(base: Path) -> Path:
    if base.exists():
        return base
    matches = glob.glob(str(base.parent.parent / "**" / base.name), recursive=True)
    return Path(matches[0]) if matches else base


def main():
    prs = build_presentation()
    slide_count = len(prs.slides)
    out_dirs = [ROOT / "engagements" / "dealality-ai-demand-positioning-deck"]
    gdrive_matches = glob.glob(r"g:/My Drive/**/Strategy & Foundations", recursive=True)
    if gdrive_matches:
        out_dirs.append(Path(gdrive_matches[0]))
    for out_dir in out_dirs:
        target = resolve_out_dir(out_dir)
        target.mkdir(parents=True, exist_ok=True)
        path = target / OUTPUT_FILENAME
        prs.save(str(path))
        print(f"Saved: {path} ({slide_count} slides)")


if __name__ == "__main__":
    main()
