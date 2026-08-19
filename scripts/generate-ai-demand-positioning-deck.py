#!/usr/bin/env python3
"""Generate Dealality AI Demand Positioning deck styled to match Short Overview Deck PDFs."""

from __future__ import annotations

import glob
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "engagements" / "dealality-ai-demand-positioning-deck" / "assets"
OUT_DIRS = [
    ROOT / "engagements" / "dealality-ai-demand-positioning-deck",
    Path(r"g:/My Drive") / "Dealality\u2122" / "Strategy & Foundations",
]

# Extracted from sample PDF typography / colors
C_BLACK = RGBColor(0x01, 0x01, 0x01)
C_GOLD = RGBColor(0xC2, 0xA8, 0x69)
C_RED = RGBColor(0xFF, 0x31, 0x31)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BODY = RGBColor(0x00, 0x00, 0x00)

FONT_HEAD = "Arial Narrow"
FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FONT_LABEL = "Arial"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
DATE_LABEL = "August 2026"


def asset(name: str) -> str:
    path = ASSETS / name
    if not path.exists():
        raise FileNotFoundError(f"Missing asset: {path}")
    return str(path)


def set_white_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_WHITE


def add_picture_cover(slide, img_path: str, left=0, top=0, width=None, height=None):
    width = width or SLIDE_W
    height = height or SLIDE_H
    return slide.shapes.add_picture(img_path, left, top, width=width, height=height)


def add_textbox(slide, left, top, width, height, text="", size=14, bold=False, color=C_BODY,
                align=PP_ALIGN.LEFT, font=FONT_BODY, wrap=True, valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_after = Pt(2)
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb, tf


def add_para(tf, text, size=14, bold=False, color=C_BODY, align=PP_ALIGN.LEFT, space_before=0, font=FONT_BODY):
    p = tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(4)
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_title_with_red_period(slide, left, top, width, title, size=52):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = title
    r1.font.name = FONT_TITLE
    r1.font.size = Pt(size)
    r1.font.bold = True
    r1.font.color.rgb = C_BLACK
    r2 = p.add_run()
    r2.text = "."
    r2.font.name = FONT_TITLE
    r2.font.size = Pt(size)
    r2.font.bold = True
    r2.font.color.rgb = C_RED
    return tb


def add_section_on_photo(slide, text, y=Inches(2.55), size=58):
    add_textbox(slide, Inches(0.45), y, Inches(12.4), Inches(0.9), text.upper(),
                size=size, bold=True, color=C_WHITE, font=FONT_HEAD)


def add_logo_corner(slide, right=True):
    icon = asset("slide03-img5-400x400.png")
    x = Inches(12.35) if right else Inches(0.45)
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(0.25), Inches(0.55), Inches(0.55))
    box.fill.solid()
    box.fill.fore_color.rgb = C_WHITE
    box.line.fill.background()
    slide.shapes.add_picture(icon, x + Inches(0.06), Inches(0.31), width=Inches(0.43))


def add_cover_logo(slide):
    # Match PDF cover: outline building mark + DEALALITY wordmark
    x0, y0 = Inches(0.55), Inches(0.45)
    for w, h, dx in [(0.18, 0.42, 0.0), (0.14, 0.32, 0.24)]:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x0 + Inches(dx), y0 + Inches(0.08), Inches(w), Inches(h))
        box.fill.background()
        box.line.color.rgb = C_BLACK
        box.line.width = Pt(2)
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x0 + Inches(0.27), y0 + Inches(0.33), Inches(0.05), Inches(0.05))
    dot.fill.solid()
    dot.fill.fore_color.rgb = C_RED
    dot.line.fill.background()
    add_textbox(slide, x0 + Inches(0.48), y0 + Inches(0.12), Inches(1.8), Inches(0.35), "DEALALITY",
                size=18, bold=True, color=C_BLACK, font=FONT_HEAD)


def add_pattern_panel(slide, left, top, width, height):
    slide.shapes.add_picture(asset("slide02-img5-2800x1575.jpeg"), left, top, width=width, height=height)


def add_vertical_banner(slide, text, x=Inches(4.95), photo_width=Inches(4.0)):
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(0), photo_width, SLIDE_H)
    banner.fill.background()
    banner.line.fill.background()
    tb = slide.shapes.add_textbox(x + Inches(0.08), Inches(1.0), Inches(0.55), Inches(5.5))
    tb.rotation = 270
    tf = tb.text_frame
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = FONT_HEAD
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = C_WHITE


def add_gold_heading_block(slide, left, top, width, title, body, title_size=21, body_size=17.5):
    add_textbox(slide, left, top, width, Inches(0.55), title, size=title_size, bold=True, color=C_GOLD, font=FONT_TITLE)
    add_textbox(slide, left, top + Inches(0.45), width, Inches(0.95), body, size=body_size, color=C_GOLD, font=FONT_BODY)


def add_hero_top_slide(slide, photo, section, headline, intro, columns, footer=""):
    add_picture_cover(slide, photo, 0, 0, SLIDE_W, Inches(3.0))
    add_logo_corner(slide)
    add_section_on_photo(slide, section, y=Inches(2.05), size=54)
    add_title_with_red_period(slide, Inches(0.55), Inches(3.15), Inches(12.2), headline, size=34)
    if intro:
        add_textbox(slide, Inches(0.55), Inches(3.95), Inches(12.2), Inches(0.55), intro, size=13, color=C_BODY)
    col_w = Inches(3.85)
    start_y = Inches(4.55 if intro else 4.05)
    for i, (title, body) in enumerate(columns[:3]):
        add_gold_heading_block(slide, Inches(0.55 + i * 4.1), start_y, col_w, title, body)
    if footer:
        add_textbox(slide, Inches(7.0), Inches(6.55), Inches(5.8), Inches(0.45), footer,
                    size=13, color=C_BODY, align=PP_ALIGN.RIGHT)


def add_split_photo_slide(slide, photo, banner_text, headline, intro, blocks, footer="", photo_width=Inches(4.35)):
    content_w = SLIDE_W - photo_width - Inches(0.15)
    add_pattern_panel(slide, Inches(0), Inches(0), content_w, SLIDE_H)
    add_picture_cover(slide, photo, content_w, 0, photo_width, SLIDE_H)
    add_vertical_banner(slide, banner_text, x=content_w - Inches(0.05), photo_width=photo_width)
    add_logo_corner(slide)
    title_size = 34 if len(headline) > 42 else 40
    add_title_with_red_period(slide, Inches(0.55), Inches(0.45), content_w - Inches(0.8), headline, size=title_size)
    y = Inches(1.35 if title_size >= 40 else 1.55)
    if intro:
        add_textbox(slide, Inches(0.55), y, content_w - Inches(0.8), Inches(0.85), intro, size=13, color=C_BODY)
        y += Inches(0.95)
    for title, body in blocks:
        add_gold_heading_block(slide, Inches(0.55), y, content_w - Inches(0.9), title, body, title_size=19, body_size=15)
        y += Inches(1.35)
    if footer:
        add_textbox(slide, Inches(0.55), Inches(6.55), content_w - Inches(0.7), Inches(0.45), footer,
                    size=13, bold=True, color=C_BLACK)


def add_cover_slide(slide):
    add_picture_cover(slide, asset("slide01-img4-1280x720.jpeg"))
    add_cover_logo(slide)
    add_textbox(slide, Inches(0.55), Inches(2.05), Inches(7.0), Inches(2.0),
                "YOUR NEXT GUEST MAY CHOOSE THEIR HOTEL BEFORE THEY EVER SEARCH FOR IT.",
                size=44, bold=True, color=C_BLACK, font=FONT_HEAD)
    add_textbox(slide, Inches(0.55), Inches(4.15), Inches(6.8), Inches(0.45),
                "INTRODUCING DEALALITY AI DEMAND POSITIONING\u2122",
                size=16, bold=True, color=C_BLACK, font=FONT_LABEL)
    add_textbox(slide, Inches(0.55), Inches(4.65), Inches(6.8), Inches(0.45),
                "FOR OWNERS: MEASURE AI CONSIDERATION BEFORE THE BOOKING BEGINS",
                size=14, bold=True, color=C_BLACK, font=FONT_LABEL)
    stat_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(5.25), Inches(6.4), Inches(1.25))
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = C_WHITE
    stat_box.line.color.rgb = C_GOLD
    stat_box.line.width = Pt(1)
    add_textbox(slide, Inches(0.8), Inches(5.35), Inches(1.5), Inches(0.7), "56%",
                size=48, bold=True, color=C_BLACK, font=FONT_TITLE)
    add_textbox(slide, Inches(2.2), Inches(5.4), Inches(4.5), Inches(0.95),
                "of U.S. travelers used AI for planning, booking or in-destination assistance in the past 12 months.\n33% \u2192 56% (+23 pts / ~70% growth)",
                size=11, color=C_BODY)
    add_textbox(slide, Inches(0.55), Inches(6.85), Inches(2.0), Inches(0.25), DATE_LABEL,
                size=14, bold=True, color=C_BLACK, font=FONT_LABEL)
    add_textbox(slide, Inches(8.8), Inches(6.85), Inches(4.0), Inches(0.25), "OWNER INTRODUCTION",
                size=12, bold=True, color=C_WHITE, font=FONT_LABEL, align=PP_ALIGN.RIGHT)


def add_close_slide(slide):
    add_picture_cover(slide, asset("slide08-img1-800x499.jpeg"))
    add_logo_corner(slide)
    add_textbox(slide, Inches(0.55), Inches(1.35), Inches(7.5), Inches(2.0),
                "DEMAND IS BEING SHAPED\nBEFORE THE BOOKING BEGINS.",
                size=48, bold=True, color=C_WHITE, font=FONT_HEAD)
    bullets = [
        "Request a private walkthrough",
        "Discuss a 3\u20135 hotel pilot",
        "Explore your portfolio use case",
    ]
    tb, tf = add_textbox(slide, Inches(0.55), Inches(3.55), Inches(5.5), Inches(1.2), bullets[0], size=16, color=C_WHITE)
    for b in bullets[1:]:
        add_para(tf, b, size=16, color=C_WHITE)
    contacts = [
        ("WEBSITE", "www.Dealality.com"),
        ("EMAIL", "joan@aohospitalityadvisors.com"),
        ("PHONE", "+34 674 993 637"),
    ]
    y = Inches(5.35)
    for label, value in contacts:
        add_textbox(slide, Inches(8.0), y, Inches(1.2), Inches(0.25), label, size=12, bold=True, color=C_WHITE, font=FONT_LABEL)
        add_textbox(slide, Inches(9.25), y, Inches(3.8), Inches(0.25), value, size=12, color=C_WHITE)
        y += Inches(0.38)
    add_textbox(slide, Inches(7.5), Inches(6.85), Inches(5.5), Inches(0.25),
                "CONFIDENTIAL INTRODUCTIONS AVAILABLE BY REQUEST.",
                size=10, bold=True, color=C_WHITE, font=FONT_LABEL, align=PP_ALIGN.RIGHT)


def add_stat_row(slide, stats, y=Inches(2.0), height=Inches(1.35)):
    width = Inches(12.2 / len(stats))
    for i, (value, label, sub) in enumerate(stats):
        x = Inches(0.55 + i * (12.2 / len(stats)))
        add_textbox(slide, x, y, width, Inches(0.65), value, size=40, bold=True, color=C_BLACK, font=FONT_TITLE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.62), width, Inches(0.45), label, size=11, bold=True, color=C_GOLD, font=FONT_TITLE, align=PP_ALIGN.CENTER)
        if sub:
            add_textbox(slide, x, y + Inches(1.0), width, Inches(0.35), sub, size=9, color=C_BODY, align=PP_ALIGN.CENTER)


def add_source_line(slide, text):
    add_textbox(slide, Inches(0.55), Inches(7.05), Inches(12), Inches(0.25), f"Source: {text}",
                size=8, color=C_GOLD, font=FONT_LABEL)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    add_cover_slide(slide)

    # 2 The Shift
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_hero_top_slide(
        slide,
        asset("slide02-img1-800x466.jpeg"),
        "THE SHIFT",
        "Hotel Discovery Is Being Rewritten",
        "The shift isn't theoretical. AI-driven travel discovery is already scaling at extraordinary rates.",
        [
            ("+2,215% AI Traffic Growth", "Growth in traffic from AI sources to U.S. travel sites from October 2024 through May 2026."),
            ("+194% YoY in May 2026", "Year-over-year growth in AI-referred travel traffic in May 2026 alone."),
            ("89% Want AI Planning", "Of consumers globally say they want to use AI in future travel planning."),
        ],
        "A new intermediary is forming between traveler demand and the hotels that get considered.",
    )
    add_textbox(slide, Inches(0.55), Inches(6.05), Inches(12), Inches(0.45),
                "SEARCH: Google \u2192 results \u2192 hotel   becomes   CONVERSATION: Traveler intent \u2192 AI interpretation \u2192 recommendation \u2192 consideration \u2192 hotel/OTA",
                size=12, bold=True, color=C_BLACK)
    add_source_line(slide, "Adobe Digital Insights, June 2026; Booking.com Global AI Sentiment Report, July 2025.")

    # 3 Why This Matters
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_hero_top_slide(
        slide,
        asset("slide06-img3-799x521.jpeg"),
        "WHY THIS MATTERS",
        "AI Is Becoming a Consideration Engine \u2014 Before It Becomes a Booking Engine",
        "AI's influence today is disproportionately concentrated in research, discovery and consideration.",
        [
            ("53% Trust AI Suggestions", "Of travelers surveyed by Expedia are comfortable allowing AI to suggest travel options."),
            ("48% Discover More", "Say AI saves time and helps them discover places they otherwise would not have found."),
            ("68% Still Book Direct", "Still prefer to ultimately book with a trusted travel brand rather than an AI chatbot or agent."),
        ],
        "Hotel owners currently have almost no performance system for measuring it.",
    )
    add_textbox(slide, Inches(0.55), Inches(5.95), Inches(12), Inches(0.45),
                "AI influences the shortlist. The hotel or OTA still captures the transaction. The emerging battleground is AI consideration.",
                size=13, bold=True, color=C_BLACK)
    add_source_line(slide, "Expedia Group AI Trust Gap, April 2026; Adobe Digital Insights, June 2026.")

    # 4 Blind Spot
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_hero_top_slide(
        slide,
        asset("slide11-img1-800x534.jpeg"),
        "THE BLIND SPOT",
        "We Measure Almost Everything \u2014 Except This",
        "Hotel owners have mature systems for measuring market performance, pricing, reputation, distribution and conversion.",
        [
            ("Market Performance", "Occupancy \u00b7 ADR \u00b7 RevPAR \u00b7 RevPAR Index"),
            ("Revenue + Digital", "Pace \u00b7 Pickup \u00b7 Segment mix \u00b7 Channel mix \u00b7 Search \u00b7 Paid media \u00b7 Conversion"),
            ("Reputation", "Ratings \u00b7 Reviews \u00b7 Guest sentiment"),
        ],
        "Dealality creates the missing performance layer.",
    )
    add_textbox(slide, Inches(0.55), Inches(5.85), Inches(12), Inches(0.55),
                "BUT\u2026 Who does AI recommend? For what demand? Against whom? And why?",
                size=20, bold=True, color=C_BLACK, font=FONT_TITLE)
    add_source_line(slide, "Deloitte, 2026 Travel Industry Outlook.")

    # 5 Product Intro (solution split layout)
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_split_photo_slide(
        slide,
        asset("slide03-img1-800x547.jpeg"),
        "THE SOLUTION",
        "Introducing AI Demand Positioning\u2122",
        "Dealality measures how a hotel competes inside AI-driven demand \u2014 not simply whether its name appears in an AI response.",
        [
            ("Are We Considered?", "How often does the hotel enter relevant AI consideration sets?"),
            ("For What Demand?", "Families? Business? Lifestyle? Groups? Experiences? Location?"),
            ("Against Whom?", "Which hotels repeatedly compete for that demand?"),
            ("Why Do We Win or Lose?", "Product? Positioning? Reviews? Narrative? Amenities? Sources?"),
        ],
        "From AI visibility to AI demand intelligence: where are we winning and losing demand \u2014 and why?",
    )

    # 6 How It Works
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_split_photo_slide(
        slide,
        asset("slide05-img5-799x449.jpeg"),
        "HOW IT WORKS",
        "An Always-On AI Demand Intelligence Engine",
        "Dealality converts thousands of individual AI interactions into structured, repeatable hotel intelligence.",
        [
            ("1. Define Demand", "Hundreds or thousands of relevant traveler intentions."),
            ("2. Observe + Capture", "Run structured demand scenarios and capture recommendations, competitors, position, narratives and sources."),
            ("3. Normalize + Benchmark", "Turn individual AI responses into comparable observations across hotel, comp set, market and portfolio."),
            ("4. Diagnose + Retest", "Identify why demand is being won or lost, then measure change over time."),
        ],
        "One prompt is anecdotal. Repeated observation becomes intelligence.",
    )

    # 7 AI Consideration Index
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_split_photo_slide(
        slide,
        asset("slide12-img1-2880x1920.jpeg"),
        "THE KPI",
        "A New Performance Metric: AI Consideration Index\u2122",
        "We want AI demand performance to be as intuitive to an owner as RevPAR Index.",
        [
            ("112 Index Score", "12% above expected competitive share in this illustrative example."),
            ("100 = Expected Share", ">100 = outperforming expected share. <100 = underperforming expected share."),
            ("Supporting KPIs", "Consideration Rate \u00b7 Competitive Win Rate \u00b7 Top Recommendation Rate \u00b7 90-Day Position Change \u00b7 Relevant Demand Coverage"),
        ],
        "Figures shown are illustrative Dealality product output, not third-party market statistics.",
    )

    # 8 Product Demo
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    content_w = SLIDE_W - Inches(4.35) - Inches(0.15)
    add_pattern_panel(slide, Inches(0), Inches(0), content_w, SLIDE_H)
    add_picture_cover(slide, asset("slide07-img1-800x533.jpeg"), content_w, 0, Inches(4.35), SLIDE_H)
    add_vertical_banner(slide, "LIVE PRODUCT", x=content_w - Inches(0.05), photo_width=Inches(4.35))
    add_logo_corner(slide)
    add_title_with_red_period(slide, Inches(0.55), Inches(0.45), content_w - Inches(0.8),
                              "AI Demand Positioning \u2014 Product View", size=34)
    add_textbox(slide, Inches(0.55), Inches(1.25), content_w - Inches(0.8), Inches(0.45),
                "Illustrative dashboard showing consideration performance, demand territories, and priority opportunities.",
                size=13, color=C_BODY)
    demo_stats = [("112", "AI Index"), ("67%", "Consideration"), ("54%", "Win Rate"), ("31%", "Family Gap"), ("+8 pts", "90-Day \u0394")]
    add_stat_row(slide, [(v, l, "") for v, l in demo_stats], y=Inches(1.85), height=Inches(1.0))
    rows = [
        ("Business Travel", "78%", "+18 pts"), ("Waterfront", "74%", "+14 pts"), ("Families", "31%", "\u221224 pts"),
    ]
    y = Inches(3.35)
    add_textbox(slide, Inches(0.55), y, Inches(4), Inches(0.25), "DEMAND POSITION MAP", size=12, bold=True, color=C_GOLD, font=FONT_TITLE)
    y += Inches(0.35)
    for terr, cons, gap in rows:
        add_textbox(slide, Inches(0.55), y, Inches(2.5), Inches(0.25), terr, size=11, color=C_BODY)
        add_textbox(slide, Inches(3.1), y, Inches(0.7), Inches(0.25), cons, size=11, bold=True, color=C_BLACK)
        add_textbox(slide, Inches(3.8), y, Inches(0.9), Inches(0.25), gap, size=11, bold=True, color=C_GOLD)
        y += Inches(0.32)
    add_textbox(slide, Inches(0.55), Inches(5.0), content_w - Inches(0.8), Inches(0.9),
                "#1 OPPORTUNITY \u2014 FAMILY DEMAND\nConsideration 31% \u00b7 Competitive gap \u221224 pts \u00b7 HIGH opportunity rating\nLimited AI evidence around room configurations, family activities, pool/recreation, and local family experiences.",
                size=12, color=C_BODY)
    add_textbox(slide, Inches(0.55), Inches(6.55), content_w - Inches(0.8), Inches(0.35),
                "From headline index to demand territories, competitors, and actionable opportunities \u2014 in one platform.",
                size=13, bold=True, color=C_BLACK)

    # 9 Observed Demand
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_split_photo_slide(
        slide,
        asset("slide04-img1-800x534.jpeg"),
        "OBSERVED DEMAND",
        "Travelers Don't Search in Hotel Segments. They Express Needs.",
        "Conversational AI allows travelers to communicate far more context than traditional keyword search.",
        [
            ("Real Prompt Examples", "\"Upscale hotel in Tampa near the waterfront for a family with two children.\" \u201cCool hotel in Waikiki that feels local rather than corporate.\""),
            ("Demand Territories", "WHO: Families \u00b7 Couples \u00b7 Business \u00b7 Groups. WHY: Business \u00b7 Weekend escape \u00b7 Event \u00b7 Family vacation. NEED: Walkability \u00b7 Dining \u00b7 Pool \u00b7 Rooms \u00b7 Design \u00b7 Location."),
            ("Booking.com Signals", "38% research destinations / timing \u00b7 37% discover local experiences \u00b7 36% find restaurant recommendations."),
        ],
        "AI gives us a new way to observe not only search volume \u2014 but demand intent.",
    )

    # 10 Demand Position Map
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_split_photo_slide(
        slide,
        asset("slide10-img2-2800x1575.jpeg"),
        "DEMAND MAP",
        "See Where Each Asset Wins \u2014 and Where It Disappears",
        "Dealality breaks overall AI performance into the demand territories that matter economically to the hotel.",
        [
            ("Strong Territories", "Business Travel 78% (+18 pts) \u00b7 Waterfront 74% (+14 pts) \u00b7 Couples 62% (+3 pts)."),
            ("Weak Territories", "Families 31% (\u221224 pts) \u00b7 Local Experience 27% (\u221229 pts)."),
            ("Family Diagnosis", "Competitor A: larger rooms. Competitor B: pool/family narrative. Competitor C: walkable family activities."),
        ],
        "This is not simply an SEO problem. It is a demand positioning problem.",
    )

    # 11 Dynamic Comp Set
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_split_photo_slide(
        slide,
        asset("slide05-img5-799x449.jpeg"),
        "COMP SET",
        "AI Does Not Know Your STR Comp Set",
        "When a traveler asks AI for a recommendation, the model constructs a competitive set dynamically around the traveler's need.",
        [
            ("Traditional Comp Set", "Hotel versus 5\u20138 predetermined competitors."),
            ("Observed Comp Set\u2122", "Hotel versus every property repeatedly entering the same AI demand conversations."),
            ("Expected / Emerging / Unexpected", "Identify expected competitors, emerging properties, unexpected takers, and demand-specific rivals."),
        ],
        "The competitive set becomes an observed outcome \u2014 not only a management assumption.",
    )

    # 12 Win / Loss
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_split_photo_slide(
        slide,
        asset("slide02-img1-800x466.jpeg"),
        "WIN / LOSS",
        "Move From \u201cWho Won?\u201d to \u201cWhy Did They Win?\u201d",
        "Dealality analyzes recurring recommendation patterns to identify the attributes, narratives and evidence driving competitive outcomes.",
        [
            ("Your Hotel Wins", "+ Location \u00b7 + Meeting capability \u00b7 + Brand awareness \u00b7 + Waterfront."),
            ("Your Hotel Loses", "\u2212 Family suitability \u00b7 \u2212 Local character \u00b7 \u2212 F&B narrative \u00b7 \u2212 Distinctive experience."),
            ("Competitor A Example", "Wins 68% of head-to-head family scenarios. Recurring reasons: room/family configuration, pool/recreation, nearby family experiences, review language."),
        ],
        "Phocuswright: AI-using travelers tend to have higher household incomes, take more trips, and spend more annually on travel.",
    )

    # 13 Narrative + Source
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_split_photo_slide(
        slide,
        asset("slide11-img1-800x534.jpeg"),
        "NARRATIVE",
        "What Does AI Believe About Your Hotel \u2014 and Why?",
        "AI recommendations synthesize information across a broad digital ecosystem \u2014 not just your own marketing message.",
        [
            ("What You Say", "Hotel website \u00b7 Brand content \u00b7 Property descriptions."),
            ("What the Market Says", "Guest reviews \u00b7 Travel publications \u00b7 Destination content \u00b7 Social conversation \u00b7 Third-party travel sites."),
            ("What AI Believes", "Who the hotel is for \u00b7 What makes it different \u00b7 Where it wins \u00b7 When it should be recommended."),
        ],
        "Hotel positioning is no longer defined only by what you publish.",
    )

    # 14 Opportunity Engine
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_split_photo_slide(
        slide,
        asset("slide03-img1-800x547.jpeg"),
        "OPPORTUNITY",
        "Don't Give Me Another Dashboard. Tell Me What to Do.",
        "The value of intelligence is identifying the actions most likely to change competitive position.",
        [
            ("#1 Family Demand", "Current consideration 31% \u00b7 Competitive gap \u221224 pts \u00b7 HIGH opportunity rating."),
            ("Why", "Limited AI evidence around room configurations, family-friendly activities, pool/recreation, and family-oriented local experiences."),
            ("What To Do", "Product \u00b7 Content \u00b7 Distribution \u00b7 Reputation \u00b7 AI Presence \u2192 then retest baseline, intervention, and observed movement."),
        ],
        "For an owner, this stops being a marketing exercise and becomes an asset-positioning tool.",
    )

    # 15 Newbond + Dovetail
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_split_photo_slide(
        slide,
        asset("slide07-img1-800x533.jpeg"),
        "USE CASE",
        "Built for the Questions Hotel Owners Actually Ask",
        "Particularly powerful across portfolios where repositioning, differentiated experiences and active asset management drive value.",
        [
            ("Newbond", "15+ hotels \u00b7 6,000+ keys \u00b7 $1.0B budgeted annual hotel revenue. Opportunistic acquisition and repositioning of upper-upscale and luxury full-service hotels."),
            ("Dovetail + Co", "Owner and creator of experiential real estate across New York, Newport, Waikiki, Bishop, and Bermuda."),
            ("Central Question", "If differentiation creates value, you should be able to measure whether the market \u2014 and AI \u2014 actually recognizes it."),
        ],
        "Significantly more valuable when monitored across an entire portfolio.",
    )

    # 16 Portfolio
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_split_photo_slide(
        slide,
        asset("slide06-img3-799x521.jpeg"),
        "PORTFOLIO",
        "One Hotel Gives You Insight. A Portfolio Gives You Intelligence.",
        "AI Demand Positioning creates a common framework for comparing positioning strength and opportunity across assets.",
        [
            ("Portfolio View", "Hotel A 118 (+9) \u00b7 Hotel B 106 (+2) \u00b7 Hotel C 97 (+7) \u00b7 Hotel D 84 (\u22126) \u00b7 Hotel E 72 (\u22124)."),
            ("Owner Use Cases", "Asset management \u00b7 Repositioning \u00b7 CapEx \u00b7 Acquisitions \u00b7 Branding \u00b7 Marketing \u00b7 Portfolio strategy."),
            ("Goal", "Establish an intelligence baseline now \u2014 while AI travel behavior is still being reshaped."),
        ],
        "Illustrative Dealality output.",
    )

    # 17 Urgency
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    add_hero_top_slide(
        slide,
        asset("slide02-img1-800x466.jpeg"),
        "WHY NOW",
        "The Window to Establish Position Is Now",
        "Traveler adoption is accelerating, AI referrals are accelerating, and AI systems are becoming more capable of recommending travel.",
        [
            ("33% \u2192 56%", "U.S. travelers using AI around travel from 1H25 to 1H26."),
            ("+2,215%", "AI-sourced U.S. travel traffic from Oct 2024 to May 2026."),
            ("3X GenAI Planning", "Increase in GenAI trip-planning adoption from 2023 to 2025."),
        ],
        "The best time to understand your AI demand position is before it becomes a board-level KPI.",
    )
    add_source_line(slide, "Phocuswright 2026; Adobe Digital Insights June 2026; Deloitte 2026 Travel Industry Outlook.")

    # 18 Pilot (beta-style)
    slide = prs.slides.add_slide(blank)
    set_white_bg(slide)
    content_w = SLIDE_W - Inches(4.35) - Inches(0.15)
    add_pattern_panel(slide, Inches(0), Inches(0), content_w, SLIDE_H)
    add_picture_cover(slide, asset("slide07-img1-800x533.jpeg"), content_w, 0, Inches(4.35), SLIDE_H)
    add_vertical_banner(slide, "PILOT", x=content_w - Inches(0.05), photo_width=Inches(4.35))
    add_logo_corner(slide)
    add_title_with_red_period(slide, Inches(0.55), Inches(0.45), content_w - Inches(0.8),
                              "Start With the Assets Where Positioning Matters Most", size=30)
    add_textbox(slide, Inches(0.55), Inches(1.25), content_w - Inches(0.8), Inches(0.45),
                "Rather than a broad commitment, Dealality can demonstrate the intelligence against a small group of real hotels.",
                size=13, color=C_BODY)
    phases = [
        ("Phase 1", "Establish demand universe"),
        ("Phase 2", "Measure baseline"),
        ("Phase 3", "Discover actual competition"),
        ("Phase 4", "Diagnose win / loss"),
        ("Phase 5", "Identify opportunities"),
        ("Phase 6", "Retest movement"),
    ]
    y = Inches(1.85)
    for i, (phase, desc) in enumerate(phases):
        col = i % 2
        row = i // 2
        x = Inches(0.55 + col * 3.8)
        yy = y + Inches(row * 0.95)
        add_gold_heading_block(slide, x, yy, Inches(3.5), phase, desc, title_size=16, body_size=13)
    add_textbox(slide, Inches(0.55), Inches(5.0), content_w - Inches(0.8), Inches(0.9),
                "Deliverables: AI Consideration Index \u00b7 Demand Position Map \u00b7 Observed Competitive Set \u00b7 Win/Loss \u00b7 Narrative Intelligence \u00b7 Source Intelligence \u00b7 Priority Opportunity Map \u00b7 Portfolio Comparison",
                size=11, color=C_BODY)
    gold = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, content_w + Inches(0.35), Inches(2.0), Inches(3.55), Inches(2.35))
    gold.fill.solid()
    gold.fill.fore_color.rgb = C_GOLD
    gold.line.fill.background()
    add_textbox(slide, content_w + Inches(0.55), Inches(2.2), Inches(3.15), Inches(0.35),
                "PROPOSED 3\u20135 HOTEL PILOT", size=13, bold=True, color=C_WHITE, font=FONT_TITLE, align=PP_ALIGN.CENTER)
    add_textbox(slide, content_w + Inches(0.55), Inches(2.75), Inches(3.15), Inches(1.35),
                "Newbond + Dovetail assets where differentiated positioning, repositioning, and portfolio intelligence matter most.",
                size=12, color=C_WHITE)
    add_textbox(slide, Inches(0.55), Inches(6.55), content_w - Inches(0.8), Inches(0.35),
                "Within one platform, ownership gets a completely new view of how its hotels compete for demand.",
                size=13, bold=True, color=C_BLACK)

    # 19 Close
    slide = prs.slides.add_slide(blank)
    add_close_slide(slide)

    return prs


def resolve_out_dir(base: Path) -> Path:
    if base.exists():
        return base
    matches = glob.glob(str(base.parent.parent / "**" / base.name), recursive=True)
    return Path(matches[0]) if matches else base


def main():
    prs = build_presentation()
    filename = "Dealality - AI Demand Positioning Deck (Owners).pptx"
    out_dirs = [ROOT / "engagements" / "dealality-ai-demand-positioning-deck"]
    gdrive = glob.glob(r"g:/My Drive/**/Strategy & Foundations", recursive=True)
    if gdrive:
        out_dirs.append(Path(gdrive[0]))
    for out_dir in out_dirs:
        target = resolve_out_dir(out_dir)
        target.mkdir(parents=True, exist_ok=True)
        path = target / filename
        prs.save(str(path))
        print(f"Saved: {path}")


if __name__ == "__main__":
    main()
